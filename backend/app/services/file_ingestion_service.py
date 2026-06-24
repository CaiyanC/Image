import hashlib
import json
from dataclasses import dataclass
from datetime import datetime, timedelta, timezone
from pathlib import Path
from typing import Any

from sqlalchemy.orm import Session

from ..models.knowledge_base import KnowledgeChunk, KnowledgeDocument


CHUNK_SIZE = 500
CHUNK_OVERLAP = 50


@dataclass
class ExtractedSection:
    text: str
    meta: dict[str, Any]


async def ingest_file(
    db: Session,
    file_path: str,
    file_name: str,
    related_skus: list[str],
    document: KnowledgeDocument | None = None,
) -> KnowledgeDocument:
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"File not found: {file_path}")

    normalized_name = (file_name or path.name or path.stem).strip()
    file_type = _resolve_file_type(path, normalized_name)
    file_hash = _calculate_file_hash(path)
    related_skus_clean = _normalize_related_skus(related_skus)

    metadata_json = json.dumps(
        {
            "file_name": normalized_name,
            "file_type": file_type,
            "file_hash": file_hash,
            "file_path": str(path),
            "related_skus": related_skus_clean,
        },
        ensure_ascii=False,
    )

    if document is None:
        document = KnowledgeDocument(
            source_type="file",
            source_id=file_hash,
            sku=related_skus_clean[0] if related_skus_clean else None,
            title=normalized_name or path.stem,
            content="",
            file_name=normalized_name,
            file_path=str(path),
            file_type=file_type,
            file_hash=file_hash,
            page_count=0,
            parse_status="processing",
            parse_error=None,
            related_skus_json=json.dumps(related_skus_clean, ensure_ascii=False),
            metadata_json=metadata_json,
        )
        db.add(document)
        db.commit()
        db.refresh(document)
    else:
        document.source_type = "file"
        document.source_id = file_hash
        document.sku = related_skus_clean[0] if related_skus_clean else document.sku
        document.title = normalized_name or path.stem
        document.content = document.content or ""
        document.file_name = normalized_name
        document.file_path = str(path)
        document.file_type = file_type
        document.file_hash = file_hash
        document.page_count = document.page_count or 0
        document.parse_status = "processing"
        document.parse_error = None
        document.related_skus_json = json.dumps(related_skus_clean, ensure_ascii=False)
        document.metadata_json = metadata_json
        db.add(document)
        db.flush()

    try:
        sections = _extract_sections(path, file_type)
        if not sections:
            raise ValueError("No text extracted from file")
        full_text = "\n\n".join(section.text for section in sections if section.text.strip()).strip()
        if not full_text:
            raise ValueError("No text extracted from file")
        chunks_payload = _build_chunks(sections, file_type=file_type)

        db.query(KnowledgeChunk).filter(KnowledgeChunk.document_id == document.id).delete(synchronize_session=False)
        document.content = full_text
        document.page_count = len(sections)
        document.parse_status = "done"
        document.parse_error = None
        document.metadata_json = json.dumps(
            {
                "file_name": normalized_name,
                "file_type": file_type,
                "file_hash": file_hash,
                "file_path": str(path),
                "page_count": len(sections),
                "chunk_count": len(chunks_payload),
                "related_skus": related_skus_clean,
            },
            ensure_ascii=False,
        )

        for item in chunks_payload:
            chunk = KnowledgeChunk(
                document_id=document.id,
                sku=related_skus_clean[0] if related_skus_clean else None,
                source_type="file",
                chunk_index=item["chunk_index"],
                content=item["content"],
                metadata_json="{}",
                embedding_status="pending",
            )
            db.add(chunk)
            db.flush()
            chunk.metadata_json = json.dumps(
                {
                    **item["metadata"],
                    "document_id": document.id,
                    "chunk_id": chunk.id,
                    "related_skus": related_skus_clean,
                },
                ensure_ascii=False,
            )
        db.commit()
        db.refresh(document)
        return document
    except Exception as exc:
        db.rollback()
        document = db.query(KnowledgeDocument).filter(KnowledgeDocument.id == document.id).first()
        if document:
            document.parse_status = "error"
            document.parse_error = str(exc)[:2000]
            document.page_count = document.page_count or 0
            document.metadata_json = json.dumps(
                {
                    "file_name": normalized_name,
                    "file_type": file_type,
                    "file_hash": file_hash,
                    "file_path": str(path),
                    "related_skus": related_skus_clean,
                },
                ensure_ascii=False,
            )
            db.commit()
            db.refresh(document)
            return document
        raise


def list_stuck_processing_documents(db: Session, timeout_minutes: int = 30) -> list[dict[str, Any]]:
    threshold = datetime.now(timezone.utc) - timedelta(minutes=max(timeout_minutes, 1))
    rows = (
        db.query(KnowledgeDocument)
        .filter(
            KnowledgeDocument.parse_status == "processing",
            KnowledgeDocument.updated_at.isnot(None),
            KnowledgeDocument.updated_at < threshold,
        )
        .all()
    )
    return [_serialize_recovery_document(document) for document in rows]


def recover_stuck_processing_documents(db: Session, timeout_minutes: int = 30) -> dict[str, Any]:
    threshold = datetime.now(timezone.utc) - timedelta(minutes=max(timeout_minutes, 1))
    rows = (
        db.query(KnowledgeDocument)
        .filter(
            KnowledgeDocument.parse_status == "processing",
            KnowledgeDocument.updated_at.isnot(None),
            KnowledgeDocument.updated_at < threshold,
        )
        .all()
    )
    recovered: list[dict[str, Any]] = []
    for document in rows:
        document.parse_status = "error"
        document.parse_error = "处理超时，可能是进程中断或解析异常"
        recovered.append(_serialize_recovery_document(document))
    if rows:
        db.commit()
        recovered = []
        for document in rows:
            db.refresh(document)
            recovered.append(_serialize_recovery_document(document))
    return {
        "recovered_count": len(rows),
        "documents": recovered,
    }


def _serialize_recovery_document(document: KnowledgeDocument) -> dict[str, Any]:
    return {
        "id": document.id,
        "file_name": document.file_name,
        "parse_status": document.parse_status,
        "updated_at": document.updated_at.isoformat() if document.updated_at else None,
        "parse_error": document.parse_error,
    }


def _resolve_file_type(path: Path, file_name: str) -> str:
    suffix = (path.suffix or Path(file_name).suffix or "").lower().lstrip(".")
    if suffix in {"pdf", "docx", "pptx", "txt", "xlsx"}:
        return suffix
    raise ValueError(f"Unsupported file type: {suffix or 'unknown'}")


def _calculate_file_hash(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as file:
        for block in iter(lambda: file.read(1024 * 1024), b""):
            digest.update(block)
    return digest.hexdigest()


def _normalize_related_skus(related_skus: list[str]) -> list[str]:
    seen: set[str] = set()
    normalized: list[str] = []
    for sku in related_skus or []:
        value = str(sku or "").strip().upper()
        if value and value not in seen:
            seen.add(value)
            normalized.append(value)
    return normalized


def _extract_sections(path: Path, file_type: str) -> list[ExtractedSection]:
    if file_type == "pdf":
        return _extract_pdf_sections(path)
    if file_type == "docx":
        return _extract_docx_sections(path)
    if file_type == "pptx":
        return _extract_pptx_sections(path)
    if file_type == "txt":
        return _extract_txt_sections(path)
    if file_type == "xlsx":
        return _extract_xlsx_sections(path)
    raise ValueError(f"Unsupported file type: {file_type}")


def _extract_pdf_sections(path: Path) -> list[ExtractedSection]:
    import fitz

    sections: list[ExtractedSection] = []
    with fitz.open(path) as doc:
        for index, page in enumerate(doc, start=1):
            text = (page.get_text("text") or "").strip()
            if not text:
                continue
            sections.append(
                ExtractedSection(
                    text=text,
                    meta={
                        "source_unit_type": "pdf_page",
                        "page_number": index,
                        "page_label": f"page-{index}",
                    },
                )
            )
    return sections


def _extract_docx_sections(path: Path) -> list[ExtractedSection]:
    from docx import Document

    document = Document(str(path))
    sections: list[ExtractedSection] = []
    for index, paragraph in enumerate(document.paragraphs, start=1):
        text = (paragraph.text or "").strip()
        if not text:
            continue
        sections.append(
            ExtractedSection(
                text=text,
                meta={
                    "source_unit_type": "docx_paragraph",
                    "paragraph_index": index,
                    "page_number": index,
                    "page_label": f"paragraph-{index}",
                },
            )
        )
    return sections


def _extract_pptx_sections(path: Path) -> list[ExtractedSection]:
    from pptx import Presentation

    presentation = Presentation(str(path))
    sections: list[ExtractedSection] = []
    for index, slide in enumerate(presentation.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.has_text_frame:
                text = (shape.text or "").strip()
                if text:
                    texts.append(text)
        combined = "\n".join(texts).strip()
        if not combined:
            continue
        sections.append(
            ExtractedSection(
                text=combined,
                meta={
                    "source_unit_type": "pptx_slide",
                    "slide_index": index,
                    "page_number": index,
                    "page_label": f"slide-{index}",
                },
            )
        )
    return sections


def _extract_txt_sections(path: Path) -> list[ExtractedSection]:
    raw = path.read_bytes()
    text, encoding = _decode_text(raw)
    cleaned = text.strip()
    if not cleaned:
        return []
    return [
        ExtractedSection(
            text=cleaned,
            meta={
                "source_unit_type": "txt_file",
                "page_number": 1,
                "page_label": "txt-1",
                "encoding": encoding,
            },
        )
    ]


def _extract_xlsx_sections(path: Path) -> list[ExtractedSection]:
    from openpyxl import load_workbook

    workbook = load_workbook(str(path), read_only=True, data_only=True)
    sections: list[ExtractedSection] = []
    try:
        for index, sheet in enumerate(workbook.worksheets, start=1):
            rows: list[str] = []
            for row in sheet.iter_rows(values_only=True):
                values = [str(cell).strip() for cell in row if cell not in (None, "")]
                if values:
                    rows.append("\t".join(values))
            text = "\n".join(rows).strip()
            if not text:
                continue
            sections.append(
                ExtractedSection(
                    text=text,
                    meta={
                        "source_unit_type": "xlsx_sheet",
                        "sheet_name": sheet.title,
                        "sheet_index": index,
                        "page_number": index,
                        "page_label": f"sheet-{index}",
                    },
                )
            )
    finally:
        workbook.close()
    return sections


def _decode_text(raw: bytes) -> tuple[str, str]:
    for encoding in ("utf-8-sig", "utf-8", "gb18030", "cp936", "big5", "latin1"):
        try:
            return raw.decode(encoding), encoding
        except UnicodeDecodeError:
            continue
    return raw.decode("utf-8", errors="ignore"), "utf-8"


def _build_chunks(sections: list[ExtractedSection], *, file_type: str) -> list[dict[str, Any]]:
    chunks: list[dict[str, Any]] = []
    chunk_index = 0
    for section_index, section in enumerate(sections, start=1):
        for item in _split_text(section.text):
            chunk_index += 1
            metadata = {
                "source_type": "file",
                "file_type": file_type,
                "section_index": section_index,
                **section.meta,
                **item["metadata"],
            }
            chunks.append(
                {
                    "chunk_index": chunk_index,
                    "content": item["content"],
                    "metadata": metadata,
                }
            )
    return chunks


def _split_text(text: str) -> list[dict[str, Any]]:
    cleaned = (text or "").strip()
    if not cleaned:
        return []
    if len(cleaned) <= CHUNK_SIZE:
        return [
            {
                "content": cleaned,
                "metadata": {
                    "start_char": 0,
                    "end_char": len(cleaned),
                    "chunk_length": len(cleaned),
                },
            }
        ]

    result: list[dict[str, Any]] = []
    start = 0
    text_length = len(cleaned)
    while start < text_length:
        end = min(start + CHUNK_SIZE, text_length)
        chunk = cleaned[start:end].strip()
        if chunk:
            result.append(
                {
                    "content": chunk,
                    "metadata": {
                        "start_char": start,
                        "end_char": end,
                        "chunk_length": len(chunk),
                    },
                }
            )
        if end >= text_length:
            break
        start = max(end - CHUNK_OVERLAP, start + 1)
    return result
