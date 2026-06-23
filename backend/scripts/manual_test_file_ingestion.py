import asyncio
import json
import sys
import tempfile
from pathlib import Path

from sqlalchemy import create_engine, select
from sqlalchemy.orm import sessionmaker

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from app.models.knowledge_base import KnowledgeChunk, KnowledgeDocument
from app.services.file_ingestion_service import ingest_file


def main() -> int:
    summary: list[dict] = []

    with tempfile.TemporaryDirectory() as temp_dir:
        temp_dir_path = Path(temp_dir)
        db_path = temp_dir_path / "test.db"

        engine = create_engine(f"sqlite+pysqlite:///{db_path}", future=True)
        SessionLocal = sessionmaker(bind=engine, autoflush=False, autocommit=False)
        try:
            KnowledgeDocument.__table__.create(bind=engine)
            KnowledgeChunk.__table__.create(bind=engine)

            with SessionLocal() as db:
                summary.append(_run_txt_case(db, temp_dir_path))
                summary.append(_run_docx_case(db, temp_dir_path))
                summary.append(_run_pptx_case(db, temp_dir_path))
                summary.append(_run_xlsx_case(db, temp_dir_path))
                summary.append(_run_pdf_case(db, temp_dir_path))
        finally:
            engine.dispose()

    print(json.dumps(summary, ensure_ascii=False, indent=2))
    return 0


def _run_txt_case(db, temp_dir_path: Path) -> dict:
    file_path = temp_dir_path / "sample.txt"
    file_path.write_text(
        "这是一个用于验证文件知识库入库流程的测试文本。\n"
        "它应该被正常解析、分块，并写入 knowledge_documents 和 knowledge_chunks。\n"
        "相关 SKU 会被序列化保存。\n"
        * 5,
        encoding="utf-8",
    )
    return _assert_case(
        db,
        file_path,
        "sample.txt",
        ["sku-001", "SKU-002", "sku-001"],
        "txt",
        expected_meta_key="page_number",
    )


def _run_docx_case(db, temp_dir_path: Path) -> dict:
    if not _can_import("docx"):
        return _skip_case("docx", "需要安装 python-docx")
    from docx import Document

    file_path = temp_dir_path / "sample.docx"
    document = Document()
    document.add_paragraph("DOCX 测试段落一。")
    document.add_paragraph("DOCX 测试段落二，包含更多文本用于分块验证。")
    document.save(str(file_path))
    return _assert_case(db, file_path, "sample.docx", ["sku-docx"], "docx", expected_meta_key="paragraph_index")


def _run_pptx_case(db, temp_dir_path: Path) -> dict:
    if not _can_import("pptx"):
        return _skip_case("pptx", "需要安装 python-pptx")
    from pptx import Presentation
    from pptx.util import Inches

    file_path = temp_dir_path / "sample.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(2))
    textbox.text_frame.text = "PPTX 测试幻灯片内容。"
    slide2 = presentation.slides.add_slide(presentation.slide_layouts[5])
    textbox2 = slide2.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(2))
    textbox2.text_frame.text = "PPTX 第二页内容，用于验证 slide 来源信息。"
    presentation.save(str(file_path))
    return _assert_case(db, file_path, "sample.pptx", ["sku-pptx"], "pptx", expected_meta_key="slide_index")


def _run_xlsx_case(db, temp_dir_path: Path) -> dict:
    if not _can_import("openpyxl"):
        return _skip_case("xlsx", "需要安装 openpyxl")
    from openpyxl import Workbook

    file_path = temp_dir_path / "sample.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "SheetOne"
    sheet.append(["名称", "值"])
    sheet.append(["字段A", "Excel 测试内容一"])
    sheet2 = workbook.create_sheet(title="SheetTwo")
    sheet2.append(["说明"])
    sheet2.append(["Excel 第二个 sheet 内容"])
    workbook.save(str(file_path))
    return _assert_case(db, file_path, "sample.xlsx", ["sku-xlsx"], "xlsx", expected_meta_key="sheet_index")


def _run_pdf_case(db, temp_dir_path: Path) -> dict:
    if not _can_import("fitz"):
        return _skip_case("pdf", "当前环境未安装 PyMuPDF，无法生成或解析 PDF")
    import fitz

    file_path = temp_dir_path / "sample.pdf"
    document = fitz.open()
    page = document.new_page()
    page.insert_text((72, 72), "PDF 测试第一页内容。")
    page2 = document.new_page()
    page2.insert_text((72, 72), "PDF 第二页内容，用于验证 page_number 来源信息。")
    document.save(str(file_path))
    document.close()
    return _assert_case(db, file_path, "sample.pdf", ["sku-pdf"], "pdf", expected_meta_key="page_number")


def _assert_case(
    db,
    file_path: Path,
    file_name: str,
    related_skus: list[str],
    file_type: str,
    *,
    expected_meta_key: str,
) -> dict:
    document = asyncio.run(
        ingest_file(
            db,
            file_path=str(file_path),
            file_name=file_name,
            related_skus=related_skus,
        )
    )

    doc = db.execute(select(KnowledgeDocument).where(KnowledgeDocument.id == document.id)).scalar_one()
    chunks = db.execute(select(KnowledgeChunk).where(KnowledgeChunk.document_id == document.id)).scalars().all()

    assert doc.source_type == "file"
    assert doc.parse_status == "done"
    assert len(chunks) >= 1
    assert all(chunk.content for chunk in chunks)
    assert all(chunk.embedding_status == "pending" for chunk in chunks)
    assert json.loads(doc.related_skus_json or "[]") == [sku.strip().upper() for sku in dict.fromkeys(related_skus)]
    assert any(expected_meta_key in json.loads(chunk.metadata_json or "{}") for chunk in chunks)

    return {
        "file_type": file_type,
        "status": "passed",
        "document_id": doc.id,
        "chunk_count": len(chunks),
        "related_skus": json.loads(doc.related_skus_json or "[]"),
    }


def _skip_case(file_type: str, reason: str) -> dict:
    return {
        "file_type": file_type,
        "status": "skipped",
        "reason": reason,
    }


def _can_import(module_name: str) -> bool:
    try:
        __import__(module_name)
        return True
    except Exception:
        return False


if __name__ == "__main__":
    raise SystemExit(main())
