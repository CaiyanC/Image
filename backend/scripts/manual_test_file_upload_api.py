import asyncio
import json
import sys
import tempfile
from pathlib import Path
from types import SimpleNamespace

from fastapi.testclient import TestClient
from sqlalchemy import create_engine, select
from sqlalchemy.orm import sessionmaker

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from app.api import knowledge_base as kb_api
from app.core.database import get_db
from app.core.security import get_current_super_admin
from app.main import app
from app.models.knowledge_base import KnowledgeChunk, KnowledgeDocument


def main() -> int:
    summary: list[dict] = []

    with tempfile.TemporaryDirectory() as temp_dir:
        temp_dir_path = Path(temp_dir)
        db_path = temp_dir_path / "test.db"
        upload_dir = temp_dir_path / "knowledge-files"
        upload_dir.mkdir(parents=True, exist_ok=True)

        engine = create_engine(f"sqlite+pysqlite:///{db_path}", future=True)
        SessionLocal = sessionmaker(bind=engine, autoflush=False, autocommit=False)
        KnowledgeDocument.__table__.create(bind=engine)
        KnowledgeChunk.__table__.create(bind=engine)

        original_upload_dir = kb_api.KNOWLEDGE_FILE_DIR
        original_startup_handlers = list(app.router.on_startup)
        original_shutdown_handlers = list(app.router.on_shutdown)
        kb_api.KNOWLEDGE_FILE_DIR = str(upload_dir)
        app.router.on_startup.clear()
        app.router.on_shutdown.clear()

        try:
            with SessionLocal() as db:
                for case in _build_positive_cases(temp_dir_path):
                    summary.append(_run_positive_case(case, db))
                summary.append(_run_duplicate_reuse_case(temp_dir_path, db))
                summary.append(_run_invalid_extension_case(temp_dir_path, db))
                summary.append(_run_oversized_file_case(temp_dir_path, db))
                summary.append(_run_path_traversal_case(temp_dir_path, db))
        finally:
            kb_api.KNOWLEDGE_FILE_DIR = original_upload_dir
            app.router.on_startup[:] = original_startup_handlers
            app.router.on_shutdown[:] = original_shutdown_handlers
            engine.dispose()

    print(json.dumps(summary, ensure_ascii=False, indent=2))
    return 0


def _run_positive_case(case: dict, db) -> dict:
    app.dependency_overrides[get_current_super_admin] = lambda: SimpleNamespace(id="super-admin")

    def override_get_db():
        try:
            yield db
        finally:
            pass

    app.dependency_overrides[get_db] = override_get_db
    try:
        response = _upload_case(case["file_path"], case["filename"], case["content_type"], case["related_skus"], db)
        assert response.status_code == 200, response.text
        payload = response.json()
        assert payload["items"], payload
        item = payload["items"][0]
        assert item["parse_status"] == "done", payload
        assert item["chunk_count"] >= 1, payload
        assert item["parse_error"] in (None, ""), payload
        assert "file_path" not in item, payload

        db.expire_all()
        doc = db.execute(select(KnowledgeDocument).where(KnowledgeDocument.id == item["document_id"])).scalar_one()
        chunks = db.execute(select(KnowledgeChunk).where(KnowledgeChunk.document_id == doc.id)).scalars().all()

        assert doc.source_type == "file"
        assert doc.parse_status == "done"
        assert json.loads(doc.related_skus_json or "[]") == [sku.strip().upper() for sku in dict.fromkeys(case["related_skus"])]
        assert len(chunks) >= 1
        assert all(chunk.content for chunk in chunks)
        assert all(chunk.embedding_status == "pending" for chunk in chunks)
        assert _has_source_marker(chunks, case["expected_meta_key"])

        return {
            "file_type": case["file_type"],
            "status": "passed",
            "document_id": doc.id,
            "chunk_count": len(chunks),
            "related_skus": json.loads(doc.related_skus_json or "[]"),
        }
    finally:
        app.dependency_overrides.pop(get_db, None)
        app.dependency_overrides.pop(get_current_super_admin, None)


def _run_invalid_extension_case(temp_dir_path: Path, db) -> dict:
    file_path = temp_dir_path / "sample.exe"
    file_path.write_text("not allowed", encoding="utf-8")
    response = _upload_case(file_path, "sample.exe", "application/octet-stream", ["SKU-ERR"], db)
    assert response.status_code == 400, response.text
    assert "不支持的文件类型" in response.text
    return {"file_type": "exe", "status": "passed", "status_code": response.status_code}


def _run_oversized_file_case(temp_dir_path: Path, db) -> dict:
    file_path = temp_dir_path / "oversized.txt"
    file_path.write_bytes(b"a" * (kb_api.MAX_KNOWLEDGE_FILE_BYTES + 1))
    response = _upload_case(file_path, "oversized.txt", "text/plain", ["SKU-ERR"], db)
    assert response.status_code == 400, response.text
    assert "文件不能超过 20MB" in response.text
    return {"file_type": "oversized", "status": "passed", "status_code": response.status_code}


def _run_path_traversal_case(temp_dir_path: Path, db) -> dict:
    file_path = temp_dir_path / "traversal.txt"
    file_path.write_text("path traversal check", encoding="utf-8")
    response = _upload_case(file_path, "../../evil.txt", "text/plain", ["SKU-TRAVERSAL"], db)
    assert response.status_code == 200, response.text
    payload = response.json()
    item = payload["items"][0]
    assert item["file_name"] == "evil.txt", payload
    return {
        "file_type": "traversal",
        "status": "passed",
        "document_id": item["document_id"],
        "chunk_count": item["chunk_count"],
        "related_skus": item["related_skus"],
    }


def _run_duplicate_reuse_case(temp_dir_path: Path, db) -> dict:
    duplicate_file = temp_dir_path / "duplicate.txt"
    duplicate_file.write_text(
        "重复文件测试内容，用于验证 file_hash 命中后复用已有文档。\n"
        "这段文本会被完整复用，不应重复生成 chunk。\n",
        encoding="utf-8",
    )

    first_response = _upload_case(duplicate_file, "duplicate.txt", "text/plain", ["SKU-DUP-A"], db)
    assert first_response.status_code == 200, first_response.text
    first_item = first_response.json()["items"][0]
    first_doc = _get_document(db, first_item["document_id"])
    first_chunk_count = _count_chunks(db, first_doc.id)

    second_response = _upload_case(duplicate_file, "duplicate.txt", "text/plain", ["SKU-DUP-B"], db)
    assert second_response.status_code == 200, second_response.text
    second_item = second_response.json()["items"][0]

    assert second_item["duplicate"] is True, second_item
    assert second_item["reused_document_id"] == first_doc.id, second_item
    assert second_item["message"] == "文件已上传过，已复用已有知识库文档", second_item
    assert second_item["chunk_count"] == first_chunk_count, second_item

    db.expire_all()
    refreshed = _get_document(db, first_doc.id)
    merged_skus = json.loads(refreshed.related_skus_json or "[]")
    assert set(merged_skus) == {"SKU-DUP-A", "SKU-DUP-B"}, merged_skus
    assert _count_chunks(db, refreshed.id) == first_chunk_count

    return {
        "file_type": "duplicate",
        "status": "passed",
        "document_id": refreshed.id,
        "chunk_count": first_chunk_count,
        "related_skus": merged_skus,
        "duplicate": True,
    }


def _upload_case(file_path: Path, filename: str, content_type: str, related_skus: list[str], db):
    app.dependency_overrides[get_current_super_admin] = lambda: SimpleNamespace(id="super-admin")

    def override_get_db():
        try:
            yield db
        finally:
            pass

    app.dependency_overrides[get_db] = override_get_db
    try:
        with TestClient(app) as client:
            return client.post(
                "/api/knowledge-base/files/upload",
                data={"related_skus": json.dumps(related_skus)},
                files=[
                    (
                        "files",
                        (
                            filename,
                            file_path.read_bytes(),
                            content_type,
                        ),
                    )
                ],
            )
    finally:
        app.dependency_overrides.pop(get_db, None)
        app.dependency_overrides.pop(get_current_super_admin, None)


def _get_document(db, document_id: str) -> KnowledgeDocument:
    return db.execute(select(KnowledgeDocument).where(KnowledgeDocument.id == document_id)).scalar_one()


def _count_chunks(db, document_id: str) -> int:
    return db.query(KnowledgeChunk).filter(KnowledgeChunk.document_id == document_id).count()


def _build_positive_cases(temp_dir_path: Path) -> list[dict]:
    return [
        _make_txt_case(temp_dir_path),
        _make_docx_case(temp_dir_path),
        _make_pptx_case(temp_dir_path),
        _make_xlsx_case(temp_dir_path),
        _make_pdf_case(temp_dir_path),
    ]


def _make_txt_case(temp_dir_path: Path) -> dict:
    file_path = temp_dir_path / "sample.txt"
    file_path.write_text(
        "这是一个用于验证上传 API 的测试文本。\n"
        "它应该被正常解析、分块，并写入 knowledge_documents 和 knowledge_chunks。\n"
        "相关 SKU 会被序列化保存。\n"
        * 5,
        encoding="utf-8",
    )
    return {
        "file_type": "txt",
        "file_path": file_path,
        "filename": file_path.name,
        "content_type": "text/plain",
        "related_skus": ["sku-001", "SKU-002", "sku-001"],
        "expected_meta_key": "page_number",
    }


def _make_docx_case(temp_dir_path: Path) -> dict:
    from docx import Document

    file_path = temp_dir_path / "sample.docx"
    document = Document()
    document.add_paragraph("DOCX 测试段落一。")
    document.add_paragraph("DOCX 测试段落二，包含更多文本用于分块验证。")
    document.save(str(file_path))
    return {
        "file_type": "docx",
        "file_path": file_path,
        "filename": file_path.name,
        "content_type": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "related_skus": ["SKU-DOCX"],
        "expected_meta_key": "paragraph_index",
    }


def _make_pptx_case(temp_dir_path: Path) -> dict:
    from pptx import Presentation
    from pptx.util import Inches

    file_path = temp_dir_path / "sample.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(2))
    textbox.text_frame.text = "PPTX 测试幻灯片内容。"
    slide2 = presentation.slides.add_slide(presentation.slide_layouts[5])
    textbox2 = slide2.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(2))
    textbox2.text_frame.text = "PPTX 第二页内容，用于验证 slide 来源信息。"
    presentation.save(str(file_path))
    return {
        "file_type": "pptx",
        "file_path": file_path,
        "filename": file_path.name,
        "content_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "related_skus": ["SKU-PPTX"],
        "expected_meta_key": "slide_index",
    }


def _make_xlsx_case(temp_dir_path: Path) -> dict:
    from openpyxl import Workbook

    file_path = temp_dir_path / "sample.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "SheetOne"
    sheet.append(["名称", "值"])
    sheet.append(["字段A", "Excel 测试内容一"])
    sheet2 = workbook.create_sheet(title="SheetTwo")
    sheet2.append(["说明"])
    sheet2.append(["Excel 第二个 sheet 内容"])
    workbook.save(str(file_path))
    return {
        "file_type": "xlsx",
        "file_path": file_path,
        "filename": file_path.name,
        "content_type": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "related_skus": ["SKU-XLSX"],
        "expected_meta_key": "sheet_index",
    }


def _make_pdf_case(temp_dir_path: Path) -> dict:
    import fitz

    file_path = temp_dir_path / "sample.pdf"
    document = fitz.open()
    page = document.new_page()
    page.insert_text((72, 72), "PDF 测试第一页内容。")
    page2 = document.new_page()
    page2.insert_text((72, 72), "PDF 第二页内容，用于验证 page_number 来源信息。")
    document.save(str(file_path))
    document.close()
    return {
        "file_type": "pdf",
        "file_path": file_path,
        "filename": file_path.name,
        "content_type": "application/pdf",
        "related_skus": ["SKU-PDF"],
        "expected_meta_key": "page_number",
    }


def _has_source_marker(chunks: list[KnowledgeChunk], expected_meta_key: str) -> bool:
    for chunk in chunks:
        try:
            metadata = json.loads(chunk.metadata_json or "{}")
        except Exception:
            metadata = {}
        if expected_meta_key in metadata:
            return True
    return False


if __name__ == "__main__":
    raise SystemExit(main())
